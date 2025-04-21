import streamlit as st
import requests
from fpdf import FPDF
from docx import Document
import openai
import PyPDF2
from docx.shared import Inches
import io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import zipfile
import os
from pathlib import Path
import csv
import pandas as pd

# Access keys from Streamlit Secrets
openai.api_key = st.secrets["api"]["OPENAI_API_KEY"]
google_api_key = st.secrets["api"]["GOOGLE_API_KEY"]
custom_search_engine_id = st.secrets["api"]["CUSTOM_SEARCH_ENGINE_ID"]



# Research-related keywords
RESEARCH_KEYWORDS = [
"research", "article", "study", "paper", "journal", "report", "document", "thesis", "dissertation",
    "review", "literature", "source", "abstract", "manuscript", "publication", "findings", "results",
    "investigation", "survey", "exploration", "clinical trial", "experiment", "data analysis",
    "methodology", "results", "hypothesis", "sample study", "case study", "data set", "research method",
    "peer-reviewed", "academic paper", "research paper", "study report", "research proposal", "field study",
    "systematic review", "experimental study", "observational study", "control group", "clinical research",
    "medical research", "pharmaceutical research", "biotech research", "drug development", "drug discovery",
    "experimental design", "epidemiological study", "randomized control trial", "meta-analysis", "biostatistics",
    "computational study", "therapeutic research", "molecular research", "genetic research", "biomedical research",
    "cancer research", "immunology study", "pathophysiology", "translational research", "treatment protocol",
    "medical innovation", "medical device study", "medical trial", "disease research", "pharmacology",
    "pharmacovigilance", "clinical development", "clinical study protocol", "patient safety", "pharmaceutical study",
    "pharmacokinetics", "therapeutic efficacy", "pharmacodynamics", "evidence-based medicine", "drug toxicology",
    "preclinical study", "clinical outcomes", "regulatory affairs", "patent study", "artificial intelligence research",
    "machine learning algorithms", "predictive modeling", "computational biology", "quantum computing research",
    "robotics in medicine", "data mining", "neural networks in pharma", "digital health", "telemedicine research",
    "precision medicine", "genomic research", "biotechnology innovation", "CRISPR technology",
    "wearable health technology",
    "peer-reviewed articles", "research journal", "scientific journal", "academic journal", "medical journal",
    "pharmaceutical journal", "research article", "review article", "open access", "editorial", "article abstract",
    "case report", "journal impact factor", "citation analysis", "scopus indexed", "elsevier", "springer",
    "wiley online library", "doi", "pubmed indexed", "neuroscience research", "cardiology research",
    "oncology research",
    "infectious disease study", "pediatrics research", "geriatrics study", "regenerative medicine",
    "stem cell research",
    "mental health studies", "HIV/AIDS research", "diabetes research", "rare diseases study",
    "autoimmune diseases research",
    "cardiovascular diseases study", "hepatology research", "dermatology study", "orthopedics research",
    "rheumatology research",
    "patent research", "patent application", "patent literature", "patent filing", "intellectual property",
    "patent search",
    "patent documentation", "pharmaceutical patent", "drug patent", "biotechnology patent", "data-driven research",
    "systematic review", "research data", "open science", "data visualization", "collaborative research",
    "research collaboration",
    "research network", "research findings", "literature review", "trial report", "cohort study",
    "cross-sectional study",
    "research grants", "clinical evaluation", "research ethics", "scientific method", "study design",
    "research funding",
    "research institutions", "research organizations", "health policy research", "epidemiology research"

]

# Function to check if a query contains research-related keywords
def is_query_research_related(query):
    for keyword in RESEARCH_KEYWORDS:
        if keyword.lower() in query.lower():
            return True
    return False

# Function to search using Google Custom Search API
def google_custom_search(query):
    url = (
        f"https://www.googleapis.com/customsearch/v1"
        f"?q={query}&key={google_api_key}&cx={custom_search_engine_id}"
    )
    response = requests.get(url)
    if response.status_code == 200:
        return response.json()
    else:
        return {"error": "Error fetching data from Google Custom Search"}

def fetch_gpt_response_content_gen(domain, query):
    try:
        system_prompt = (
            f"You are an expert in the {domain} domain only. "
            f"Only answer the questions related to the specified {domain} domain "
            "and don't answer any other questions."
        )
        
        response = openai.chat.completions.create(
            model="gpt-3.5-turbo",
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": query},
            ],
        )
        return response.choices[0].message.content
    except Exception as e:
        return f"Error: {str(e)}"




# Function to extract text from PDF
def extract_text_from_pdf(pdf_file):
    pdf_reader = PyPDF2.PdfReader(pdf_file)
    pdf_text = ""
    for page in pdf_reader.pages:
        pdf_text += page.extract_text()
    return pdf_text


def save_as_scorm_pdf(content, output_folder="scorm_package", scorm_zip_name="scorm_package.zip"):
    # Step 1: Create the SCORM folder structure
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)

    # Save the PDF
    pdf_file_path = os.path.join(output_folder, "content.pdf")
    save_as_pdf(content, pdf_file_path)

    # Step 2: Create the HTML file
    html_file_path = os.path.join(output_folder, "index.html")
    with open(html_file_path, "w", encoding="utf-8") as html_file:
        html_file.write(f"""
        <!DOCTYPE html>
        <html>
        <head>
            <title>SCORM Content</title>
        </head>
        <body>
            <h1>Research Content Response</h1>
            <iframe src="content.pdf" width="100%" height="600px"></iframe>
        </body>
        </html>
        """)

    # Step 3: Create the imsmanifest.xml file
    manifest_file_path = os.path.join(output_folder, "imsmanifest.xml")
    with open(manifest_file_path, "w", encoding="utf-8") as manifest_file:
        manifest_file.write(f"""
        <?xml version="1.0" encoding="UTF-8"?>
        <manifest xmlns="http://www.imsglobal.org/xsd/imscp_v1p1"
                  xmlns:adlcp="http://www.adlnet.org/xsd/adlcp_v1p3"
                  xmlns:xsi="http://www.w3.org/2001/XMLSchema-instance"
                  xsi:schemaLocation="http://www.imsglobal.org/xsd/imscp_v1p1
                                      http://www.imsglobal.org/xsd/imscp_v1p1.xsd
                                      http://www.adlnet.org/xsd/adlcp_v1p3
                                      http://www.adlnet.org/xsd/adlcp_v1p3.xsd">
            <metadata>
                <schema>ADL SCORM</schema>
                <schemaversion>1.2</schemaversion>
            </metadata>
            <organizations>
                <organization identifier="ORG-1">
                    <title>Research Content</title>
                    <item identifier="ITEM-1" identifierref="RES-1">
                        <title>Research Content Response</title>
                    </item>
                </organization>
            </organizations>
            <resources>
                <resource identifier="RES-1" type="webcontent" href="index.html">
                    <file href="index.html"/>
                    <file href="content.pdf"/>
                </resource>
            </resources>
        </manifest>
        """)

    # Step 4: Zip the SCORM package
    with zipfile.ZipFile(scorm_zip_name, 'w', zipfile.ZIP_DEFLATED) as scorm_zip:
        for foldername, subfolders, filenames in os.walk(output_folder):
            for filename in filenames:
                file_path = os.path.join(foldername, filename)
                arcname = os.path.relpath(file_path, output_folder)
                scorm_zip.write(file_path, arcname)

    # Provide the download button for the SCORM package
    with open(scorm_zip_name, "rb") as scorm_file:
        st.download_button("Download SCORM Package", scorm_file, scorm_zip_name, "application/zip")


def save_as_pdf(content, file_name="response.pdf"):
    pdf = FPDF()
    pdf.add_page()

    # Add the logo
    pdf.image('assets/logo.jpeg', x=10, y=8, w=30)

    # Title of the document
    pdf.set_font("Arial", style='B', size=16)
    pdf.ln(30)
    pdf.cell(200, 10, txt="Research Content Response", ln=True, align='C')
    pdf.ln(10)

    # Add content
    pdf.set_font("Arial", size=12)
    pdf.multi_cell(190, 10, content)

    # Save the PDF
    pdf.output(file_name)


def save_as_scorm_word(content, file_name="scorm_package.zip"):
    # Create an in-memory zip file
    scorm_zip = io.BytesIO()

    with zipfile.ZipFile(scorm_zip, 'w') as zf:
        # Create and add manifest.xml
        manifest_content = """<manifest>
            <metadata>
                <schema>ADL SCORM</schema>
                <schemaversion>1.2</schemaversion>
            </metadata>
            <resources>
                <resource identifier="res1" type="webcontent" href="response.docx">
                    <file href="response.docx"/>
                    <file href="response.html"/>
                </resource>
            </resources>
        </manifest>"""
        zf.writestr("imanifest.xml", manifest_content)

        # Create DOCX file
        docx_buffer = io.BytesIO()
        doc = Document()
        # Add the logo to the Word document
        logo_path = "assets/logo.jpeg"
        if Path(logo_path).is_file():
            doc.add_picture(logo_path, width=Inches(1.5))
        doc.add_paragraph('\n')
        doc.add_paragraph("Research Content Response", style='Heading 1')
        doc.add_paragraph('\n')
        doc.add_paragraph(content)
        doc.save(docx_buffer)
        docx_buffer.seek(0)
        zf.writestr("response.docx", docx_buffer.getvalue())

        # Create HTML file
        html_content = f"""
        <html>
        <head><title>Research Content Response</title></head>
        <body>
        <h1>Research Content Response</h1>
        <p>{content.replace('\n', '<br>')}</p>
        </body>
        </html>
        """
        zf.writestr("index.html", html_content)

    scorm_zip.seek(0)
    return scorm_zip.getvalue()


# Usage in Streamlit
def save_as_scorm_button(content):
    scorm_data = save_as_scorm_word(content)
    st.download_button(
        label="Download SCORM Package",
        data=scorm_data,
        file_name="scorm_package.zip",
        mime="application/zip"
    )

def fetch_gpt_response_pdf(query):
    try:
        response = openai.chat.completions.create(
            model="gpt-3.5-turbo",
            messages=[
                {"role": "system", "content": "You are an expert in analyzing PDFs and providing highlights, summaries, analyses, and insights. Only answer questions based strictly on the content of the uploaded PDF. Do not answer any questions that are unrelated or outside the scope of the PDF."},
                {"role": "user", "content": query},
            ],
        )
        return response.choices[0].message.content
    except Exception as e:
        return f"Error: {str(e)}"



def get_response(text):
    try:
        response = openai.chat.completions.create(
            model="gpt-3.5-turbo",
            messages=[{"role": "user", "content": text}]
        )
        return response.choices[0].message.content
    except Exception as e:
        return f"Error: {str(e)}"


# Function to fetch medical and pharma-related data from GPT-3
# Function to fetch data from GPT-3 based on domain and query

# Function to fetch structured CSV data from GPT-3
def fetch_gpt_response(domain, query):
    prompt = f"""
    Please provide reliable and accurate data related to the following query in the domain of {domain}.
    Don't answer queries or provide CSV data for any other domain except the one provided by the user.
    The data should include at least 15 to 20 entries and be formatted as a proper CSV with headers and rows.
    
    The response **must** strictly follow this format:
    
    ```
    Column1,Column2,Column3
    Value1,Value2,Value3
    Value4,Value5,Value6
    ```
    
    Ensure that the output is structured properly as CSV without additional text, explanations, or formatting.
    
    Query: {query}
    """
    response = get_response(prompt)  # Fetch response from GPT-3
    return response.strip()

# Function to create SCORM package dynamically based on domain and query
def create_scorm_package(csv_content, domain, query):
    # Create an in-memory binary stream for the zip file
    zip_buffer = io.BytesIO()

    with zipfile.ZipFile(zip_buffer, "w") as zip_file:
        # Add the CSV content to the zip file
        zip_file.writestr("data.csv", csv_content)

        # Dynamically create imsmanifest.xml content
        imsmanifest_content = f"""<?xml version="1.0" encoding="UTF-8"?>
<manifest identifier="scorm_2004" version="1.0">
    <organizations>
        <organization identifier="org_1">
            <title>{domain} SCORM Package</title>
        </organization>
    </organizations>
    <resources>
        <resource identifier="res_1" type="webcontent" href="index.html">
            <file href="data.csv"/>
            <file href="index.html"/>
        </resource>
    </resources>
</manifest>"""
        zip_file.writestr("imsmanifest.xml", imsmanifest_content)

        # Create dynamic index.html content with the domain and query
        index_html_content = f"""<!DOCTYPE html>
<html>
<head>
    <title>{domain} Data</title>
</head>
<body>
    <h1>Welcome to the {domain} SCORM Package</h1>
    <p><strong>Query:</strong> {query}</p>
    <p>This package contains generated data based on the domain and query provided.</p>
</body>
</html>
"""
        zip_file.writestr("index.html", index_html_content)

    # Rewind the buffer to the beginning
    zip_buffer.seek(0)
    return zip_buffer







# Function to convert CSV string to DataFrame
def csv_to_dataframe(csv_string):
    try:
        df = pd.read_csv(io.StringIO(csv_string))
        return df
    except Exception as e:
        return None  # Handle invalid CSV cases

def generate_detailed_ppt_content(domain, topic):
    """Generate detailed content for a presentation using GPT based on the selected domain and topic."""
    prompt = (
        f"You are an expert in the {domain} domain. Generate a professional, formal PowerPoint presentation on the topic: '{topic}'.\n\n"
        f"Instructions:\n"
        f"1. All content must be specific to the {domain} domain and based on the topic '{topic}'.\n"
        f"2. Each slide must have at least 4 well-written bullet points, not paragraphs.\n"
        f"3. Ensure formal tone, clarity, and relevance to the chosen domain.\n"
        f"4. Structure should include:\n"
        f"   - Title Slide (Domain, Topic, Author Name placeholder)\n"
        f"   - Introduction Slide (Definition and importance of the topic)\n"
        f"   - 4‚Äì6 Key Point Slides (with elaborated points)\n"
        f"   - Case Studies/Examples Slide (with real-world relevance to the domain)\n"
        f"   - Conclusion Slide (with summary/future direction)\n\n"
        f"Output only detailed slide-wise content."
    )
    try:
        response = openai.chat.completions.create(
            model="gpt-4",
            messages=[
                {"role": "system", "content": f"You are a domain expert in {domain}."},
                {"role": "user", "content": prompt},
            ],
        )
        return response.choices[0].message.content
    except Exception as e:
        return f"Error: {str(e)}"

def create_professional_ppt(content, topic, file_name="presentation.pptx"):
    """Create a well-formatted professional PowerPoint presentation."""
    ppt = Presentation()

    # Set consistent font styles
    def set_textbox_style(text_frame):
        """Style the textbox content."""
        for paragraph in text_frame.paragraphs:
            paragraph.font.name = "Calibri (Body)"
            paragraph.font.size = Pt(20)
            paragraph.alignment = PP_ALIGN.LEFT

    # Title Slide
    title_slide_layout = ppt.slide_layouts[0]
    slide = ppt.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = topic
    subtitle.text = f"A Comprehensive Overview in {domain} Domain"

    # Process GPT content into slides
    sections = content.split("\n\n")
    for section in sections:
        if ":" in section:  # Detect title: content structure
            slide_title, slide_content = section.split(":", 1)

            # Add a new slide for each section
            slide = ppt.slides.add_slide(ppt.slide_layouts[1])
            title = slide.shapes.title
            title.text = slide_title.strip()

            # Add content to the slide
            content_box = slide.placeholders[1]
            content_box.text = slide_content.strip()
            set_textbox_style(content_box.text_frame)

    # Save the presentation
    ppt.save(file_name)
    return file_name

# Streamlit UI
st.set_page_config(page_title="Content and Research Analysis", layout="wide")
st.title("üìö Content Generation And Analysis System")

# Sidebar Navigation
sections = ["About", "Content Generation", "PDF Analysis","CSV Content Generation", "Research Search","PPT Development", "Instructions", "Credits"]
selected_section = st.sidebar.selectbox("Navigation", sections)

if selected_section == "About":
    st.markdown("---")
    st.header("üìñ About")

    st.info("### 1. What is this application about?")
    st.success(
        "This application integrates OpenAI's GPT-3.5 and Google Custom Search API to provide a seamless platform for "
        "content generation, PDF analysis, CSV data creation, research exploration, and PowerPoint presentation development. "
        "It is designed to cater to various domains, making it a versatile tool for professionals, researchers, and educators."
    )

    st.info("### 2. What does the Content Generation section offer?")
    st.success(
        "Users can input queries related to any field to generate detailed content using GPT-3.5. "
        "The content can be downloaded as SCORM packages in PDF or Word formats."
    )

    st.info("### 3. How does the PDF Analysis section work?")
    st.success(
        "Upload a PDF document, extract text from it, and ask questions based on the content. Responses are generated using "
        "GPT-3.5 and can be downloaded as SCORM-compatible files."
    )

    st.info("### 4. What can I do in the CSV Content Generation section?")
    st.success(
        "Generate structured CSV data for different domains based on user queries. The CSV files can be packaged and downloaded as SCORM-compatible files."
    )

    st.info("### 5. What features are available in the Research Search section?")
    st.success(
        "Input a research query to fetch results from trusted sources using the Google Custom Search API. The system ensures that results are relevant and reliable."
    )

    st.info("### 6. How does the PPT Development section help?")
    st.success(
        "Enter a topic to generate a professionally crafted PowerPoint presentation. Download the PPT directly after content generation."
    )

    st.info("### 7. Who should use this application?")
    st.success(
        "This tool is ideal for professionals, researchers, educators, and content creators who require tailored content generation, "
        "research insights, and analysis tools across various domains."
    )
    # Horizontal line
    st.markdown("---")

    # Footer
    st.caption("Developed by **Corbin Technology Solutions**")

    


elif selected_section == "Content Generation":
    st.markdown("---")
    st.header("üîç Content Generation")
    
    # User selects the domain first
    domain = st.text_input(
        "Enter the domain in which the answer is required:",
        placeholder="Example: Medical, Pharmaceutical, Finance, etc."
    )

    # Ensure session state exists for response storage
    if "generated_response" not in st.session_state:
        st.session_state.generated_response = None

    if domain:
        query = st.text_area(
            "Enter your query below:",
            height=200,
            placeholder=f"Enter any query related to the {domain} domain",
        )

        if query:
            # Check if a new query has been entered
            if query != st.session_state.get("last_query"):
                # Fetch response and store in session state
                st.session_state.generated_response = fetch_gpt_response_content_gen(domain, query)
                st.session_state.last_query = query  # Update last query

            # Display the response
            st.subheader("Response")
            st.write(st.session_state.generated_response)

        # Horizontal line before download options
        st.markdown("---")

        # Download options
        st.subheader("üì• Download Options")

        # Button to download SCORM PDF
        if st.button("Generate the PDF as SCORM Package"):
            save_as_scorm_pdf(st.session_state.generated_response)
            st.success("SCORM package generated successfully!")

        # Button to download SCORM Word
        if st.button("Generate the Word File as SCORM Package"):
            scorm_word = save_as_scorm_word(st.session_state.generated_response, file_name="response.docx")
            if scorm_word:
                st.success("SCORM Word package generated successfully!")
                st.download_button(
                    label="Download SCORM Word Package",
                    data=scorm_word,
                    file_name="scorm_word_package.zip",
                    mime="application/zip",
                )
            else:
                st.error("Failed to generate SCORM Word package.")

    # Horizontal line
    st.markdown("---")

    # Footer
    st.caption("Developed by **Corbin Technology Solutions**")

elif selected_section == "PDF Analysis":
    
    # Horizontal line
    st.markdown("---")
    st.header("üìÑ PDF Analysis")

    # Upload PDF
    pdf_file = st.file_uploader("Upload a PDF", type="pdf")

    # Initialize variables
    extracted_text = ""
    response = ""

    if pdf_file:
        with io.BytesIO(pdf_file.read()) as pdf_stream:
            extracted_text = extract_text_from_pdf(pdf_stream)

        # Show Extracted Text
        st.write("Extracted Text:")
        st.text_area("PDF Content", extracted_text, height=200)

        # Ask a question based on the PDF
        query = st.text_input("Ask a question based on the PDF:")
        
        if query:
            # Only answer questions related to the uploaded PDF
            full_prompt = (
                "You are an expert in analyzing PDFs and providing highlights, summaries, analyses, and insights. "
                "Only answer questions based strictly on the content of the uploaded PDF. "
                "Do not answer any questions that are unrelated or outside the scope of the PDF.\n\n"
                f"PDF Content: {extracted_text}\n\n"
                f"Question: {query}"
            )

            # Get response
            response = fetch_gpt_response_pdf(full_prompt)

            # Display the response
            st.subheader("Response")
            st.write(response)

            # Horizontal line
            st.markdown("---")

            # Download Options
            st.subheader("Download Options")

            if st.button("Generate the Response as PDF SCORM Package"):
                save_as_scorm_pdf(response)
                st.success("SCORM package generated. Check the 'Download SCORM Package' button.")

            if st.button("Generate the Response as Word SCORM Package"):
                scorm_word = save_as_scorm_word(response, file_name="response.docx")

                if scorm_word:
                    st.success("SCORM package generated. Click the 'Download SCORM Package' button below.")
                    st.download_button(
                        label="Download SCORM Word Package",
                        data=scorm_word,
                        file_name="scorm_word_package.zip",
                        mime="application/zip"
                    )
                else:
                    st.error("Failed to generate SCORM Word package.")

    else:
        st.info("Please upload a PDF file to begin analysis.")

    # Horizontal line
    st.markdown("---")

    # Footer
    st.caption("Developed by **Corbin Technology Solutions**")
    
# Streamlit Section to Handle User Input and SCORM Package Generation
elif selected_section == "CSV Content Generation":
    # Horizontal line
    st.markdown("---")

    st.header("üîç CSV Content Generation")

    # User selects the domain first
    domain = st.text_input(
        "Enter the domain in which the answer is required:",
        placeholder="Example: Medical, Pharmaceutical, Finance, etc."
    )

    # Ensure session state exists for response storage
    if "generated_response" not in st.session_state:
        st.session_state.generated_response = None

    if domain:
        query = st.text_area(
            "Enter your query below:",
            height=200,
            placeholder=f"Enter any query related to the {domain} domain",
        )

        if query:
            # Check if a new query has been entered
            if query != st.session_state.get("last_query"):
                # Fetch response and store in session state
                st.session_state.generated_response = fetch_gpt_response(domain, query)
                st.session_state.last_query = query  # Update last query

            # Convert response to CSV format and display
            csv_data = st.session_state.generated_response
            df = csv_to_dataframe(csv_data)

            if df is not None:
                st.subheader("CSV Data Preview")
                st.dataframe(df)  # Display CSV as table

                # Horizontal line before download options
                st.markdown("---")

                # Provide a button to download the CSV file
                csv_buffer = io.StringIO()
                df.to_csv(csv_buffer, index=False)
                st.download_button(
                    label="Download CSV File",
                    data=csv_buffer.getvalue(),
                    file_name=f"{domain.lower().replace(' ', '_')}_data.csv",
                    mime="text/csv"
                )

                # Button to generate and download the CSV as a SCORM package
                if st.button("Generate SCORM Package"):
                    scorm_package = create_scorm_package(csv_data, domain, query)
                    st.download_button(
                        label="Download CSV File as SCORM Package",
                        data=scorm_package.getvalue(),
                        file_name=f"{domain.lower().replace(' ', '_')}_scorm.zip",
                        mime="application/zip"
                    )
            else:
                st.warning("‚ö† The generated response is not in a valid CSV format.")

    # Horizontal line
    st.markdown("---")

    # Footer
    st.caption("Developed by **Corbin Technology Solutions**")

# Research Search Section
elif selected_section == "Research Search":
    st.markdown("---")
    st.header("üî¨ Research Search")

    query = st.text_area(
    "Enter a research query below:",
    height=200,
    placeholder="Example: cancer treatment, new medicines, vaccine updates"
    )

    if query:
        search_results = google_custom_search(query)
        relevant_content = []
        for item in search_results.get("items", []):
            title = item.get("title", "")
            snippet = item.get("snippet", "")
            if any(keyword in title.lower() or keyword in snippet.lower() for keyword in RESEARCH_KEYWORDS):
                relevant_content.append({"title": title, "link": item.get("link", ""), "snippet": snippet})
        if relevant_content:
            st.write("**Research Results:**")
            for content in relevant_content:
                st.write(f"- **[{content['title']}]({content['link']})**")
                st.write(content["snippet"])
        else:
            st.write("No relevant research-related content found.")
     # Horizontal line
    st.markdown("---")

    # Footer
    st.caption("Developed by **Corbin Technology Solutions**")


# Streamlit Integration for PPT Generation
elif selected_section == "PPT Development":
    st.markdown("---")
    st.header("üìä PPT Content Generation")

    # Step 1: Get the domain
    domain = st.text_input(
        "Enter the domain for your presentation:", 
        placeholder="e.g., Medical, Finance, Education"
    )

    # Step 2: Get the topic
    topic = ""
    if domain:
        topic = st.text_input(
            f"Enter the topic related to the {domain} domain:",
            placeholder="e.g., Drug Discovery, Stock Market Trends, Online Learning Platforms"
        )

    # Step 3: Generate and preview content before PPT download
    if st.button("Generate PPT"):
        if domain and topic:
            st.info("Generating detailed content for your presentation. Please wait...")
            detailed_content = generate_detailed_ppt_content(domain, topic)

            if "Error" not in detailed_content:
                st.success("PowerPoint presentation generated successfully!")

                # Display structured content preview
                st.subheader("üìÑ Generated Content Preview")
                for i, slide in enumerate(detailed_content, 1):
                    st.write(f"**Slide {i}: {slide['title']}**")
                    st.write(f"{slide['content']}\n")

                # Show slide count
                st.write(f"üî¢ **Total Slides:** {slide_count}")

                # Download button
                st.download_button(
                    "üì• Download Your PPT",
                    ppt_bytes,
                    file_name=f"{domain}_{topic}.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                )
            else:
                st.error(detailed_content)
        else:
            st.warning("Please enter both domain and topic before generating the presentation.")

    # Horizontal line
    st.markdown("---")

    # Footer
    st.caption("Developed by **Corbin Technology Solutions**")


elif selected_section == "Instructions":
    st.markdown("---")
    st.header("üìù Instructions")
    st.info("Follow the steps below to effectively use each feature in the application.")

    # General Information
    with st.expander("üìã General Usage Guidelines"):
        st.markdown("""
        - This is a **generalized system** designed for multiple domains (e.g., education, technology, business, etc.).
        - Ensure your queries are relevant, clear, and domain-specific for the best results.
        - Upload files (PDFs, CSVs) in supported formats.
        - Always review generated content before downloading or sharing.
        """)

    st.subheader("üîç Feature-specific Instructions")

    # Content Generation Instructions
    with st.expander("1Ô∏è‚É£ **Content Generation**"):
        st.markdown("""
        - Use this module to generate content across different domains.
        - **Steps**:
          1. Enter the **domain** in which the query is to be asked.
          2. Enter your **query** in the text area.
          3. Click the **Submit** button to get a detailed response.
          4. Download the response as a **PDF SCORM** or **Word SCORM Package**.
        """)

    # PDF Analysis Instructions
    with st.expander("2Ô∏è‚É£ **PDF Analysis**"):
        st.markdown("""
        - Analyze and extract content from uploaded PDFs and ask domain-specific questions.
        - **Steps**:
          1. Upload a **PDF** file.
          2. Ask your **query** based on the content.
          3. Click **Submit** to obtain a relevant response.
          4. Download the response as a **PDF SCORM** or **Word SCORM Package**.
        """)

    # CSV Content Generation Instructions
    with st.expander("3Ô∏è‚É£ **CSV Content Generation**"):
        st.markdown("""
        - Generate structured CSV data based on your input query across various domains.
        - **Steps**:
          1. Enter the **domain** for which the query is to be asked.
          2. Enter your **query**.
          3. Click **Submit** to generate CSV content.
          4. Download as a **CSV file** or **CSV SCORM Package**.
        """)

    # Research Search Instructions
    with st.expander("4Ô∏è‚É£ **Research Search**"):
        st.markdown("""
        - Retrieve research papers, articles, journals, and other relevant resources.
        - **Steps**:
          1. Enter your **research query**.
          2. Click **Submit** to get search results.
          3. Access links to relevant research material directly from the results.
        """)

    # PPT Development Instructions
    with st.expander("5Ô∏è‚É£ **PPT Development**"):
        st.markdown("""
        - Create high-quality PowerPoint presentations based on any topic or domain.
        - **Steps**:
          1. Enter the **domain** and your **query/topic**.
          2. Click **Submit** to generate a presentation.
          3. Download the generated **PPT** file using the download button.
        """)

    # Footer Success Message
    st.success("Refer to these instructions for smooth navigation and utilization of all features!")
    # Horizontal line
    st.markdown("---")

    # Footer
    st.caption("Developed by **Corbin Technology Solutions**")


elif selected_section == "Credits":
    st.markdown("---")
    
    st.header("üë®‚Äçüíª Credits")

    # Highlight the developer information
    st.subheader("üåü Developed By")
    st.markdown("""
    **Corbin Technology Solutions**  
    Bringing innovative solutions for AI Chatbots and LMS Systems.
    """)

    # Technologies used
    st.subheader("üõ†Ô∏è Technologies Used")
    st.markdown("""
    - **OpenAI GPT**: For intelligent and context-aware content generation.
    - **Google Custom Search API**: For fetching trusted research content.
    - **Streamlit**: For building a modern and interactive web interface.
    """)

    # Acknowledgment
    st.subheader("üôè Acknowledgment")
    st.info("""
    Special thanks to the team at **Corbin Technology Solutions** for their dedication and expertise in creating this application.
    """)

    # Footer Message
    st.success("We appreciate your support and feedback to enhance this application!")
